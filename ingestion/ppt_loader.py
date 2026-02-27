from pptx import Presentation
from typing import List, Dict


def load_ppt(file_path: str) -> List[Dict]:
    """
    Extract text from PPT/PPTX slide by slide.
    Returns structured content with metadata.
    """

    documents = []

    try:
        presentation = Presentation(file_path)

        for slide_number, slide in enumerate(presentation.slides):
            slide_text = ""

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"

            if slide_text.strip():
                documents.append({
                    "content": slide_text,
                    "metadata": {
                        "source": file_path,
                        "slide": slide_number + 1,
                        "type": "ppt"
                    }
                })

    except Exception as e:
        raise RuntimeError(f"Error loading PPT: {str(e)}")

    return documents